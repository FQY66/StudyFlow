"""StudyFlow 毕业答辩PPT生成脚本 - 重制版"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os, copy

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── 配色方案 ──
C_PRIMARY = RGBColor(0x1A, 0x3C, 0x6E)  # 深蓝
C_SECONDARY = RGBColor(0x2B, 0x5C, 0xA8)  # 主蓝
C_ACCENT = RGBColor(0xE8, 0x6C, 0x00)  # 橙色
C_LIGHT_BG = RGBColor(0xE8, 0xF0, 0xFE)  # 浅蓝背景
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK = RGBColor(0x33, 0x33, 0x33)
C_MED = RGBColor(0x66, 0x66, 0x66)
C_LIGHT = RGBColor(0x99, 0x99, 0x99)
C_GREEN = RGBColor(0x2E, 0x7D, 0x32)
C_TEAL = RGBColor(0x00, 0x80, 0x80)
C_PURPLE = RGBColor(0x6A, 0x1B, 0x9A)


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color, radius=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        left,
        top,
        width,
        height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if radius:
        shape.adjustments[0] = radius
    return shape


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    size=18,
    bold=False,
    color=C_DARK,
    align=PP_ALIGN.LEFT,
    font="微软雅黑",
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font
    p.alignment = align
    return tb


def add_multi_text(
    slide,
    left,
    top,
    width,
    height,
    lines,
    size=16,
    color=C_DARK,
    spacing=1.4,
    font="微软雅黑",
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = font
        p.space_after = Pt(size * (spacing - 1))
    return tb


def add_section_bar(slide, title, subtitle=None):
    """统一的章节标题栏"""
    add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.07), C_SECONDARY)
    # 左侧色块
    add_rect(slide, Inches(0), Inches(0.2), Inches(0.12), Inches(0.6), C_ACCENT)
    add_textbox(
        slide,
        Inches(0.5),
        Inches(0.25),
        Inches(11),
        Inches(0.6),
        title,
        size=30,
        bold=True,
        color=C_PRIMARY,
    )
    if subtitle:
        add_textbox(
            slide,
            Inches(0.5),
            Inches(0.85),
            Inches(11),
            Inches(0.35),
            subtitle,
            size=15,
            color=C_MED,
        )


def add_card(
    slide,
    left,
    top,
    width,
    height,
    title,
    items,
    title_color=C_SECONDARY,
    bg=C_LIGHT_BG,
):
    add_rect(slide, left, top, width, height, bg, 0.05)
    add_textbox(
        slide,
        left + Inches(0.2),
        top + Inches(0.12),
        width - Inches(0.4),
        Inches(0.35),
        title,
        size=17,
        bold=True,
        color=title_color,
    )
    y = top + Inches(0.55)
    for item in items:
        add_textbox(
            slide,
            left + Inches(0.35),
            y,
            width - Inches(0.55),
            Inches(0.3),
            f"• {item}",
            size=13,
            color=C_DARK,
        )
        y += Inches(0.32)


# ================================================================
# 第1页：封面（重制版）
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_PRIMARY)

# 顶部装饰带
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.12), C_ACCENT)
# 左侧竖条
add_rect(
    slide,
    Inches(0),
    Inches(0),
    Inches(0.2),
    prs.slide_height,
    RGBColor(0x15, 0x30, 0x5A),
)
# 底部装饰带
add_rect(
    slide,
    Inches(0),
    Inches(7.0),
    prs.slide_width,
    Inches(0.5),
    RGBColor(0x15, 0x30, 0x5A),
)

# 中间装饰线
add_rect(slide, Inches(0.8), Inches(3.5), Inches(11.7), Inches(0.03), C_ACCENT)

# 学校名称
add_textbox(
    slide,
    Inches(0.8),
    Inches(0.8),
    Inches(11.7),
    Inches(0.6),
    "桂林电子科技大学",
    size=24,
    color=RGBColor(0xAA, 0xBB, 0xDD),
    align=PP_ALIGN.CENTER,
)
add_textbox(
    slide,
    Inches(0.8),
    Inches(1.3),
    Inches(11.7),
    Inches(0.4),
    "毕业设计答辩",
    size=18,
    color=RGBColor(0x88, 0xAA, 0xCC),
    align=PP_ALIGN.CENTER,
)

# 主标题
add_textbox(
    slide,
    Inches(1.5),
    Inches(2.0),
    Inches(10.3),
    Inches(0.8),
    "StudyFlow 高校研学交流平台",
    size=46,
    bold=True,
    color=C_WHITE,
    align=PP_ALIGN.CENTER,
)
add_textbox(
    slide,
    Inches(1.5),
    Inches(2.8),
    Inches(10.3),
    Inches(0.5),
    "的设计与实现",
    size=26,
    color=RGBColor(0xCC, 0xDD, 0xEE),
    align=PP_ALIGN.CENTER,
)

# 答辩信息
info_lines = [
    "答辩人：李硕",
    "学    号：XXXXXXXXXX",
    "班    级：XXX",
    "指导老师：王慧娇",
]
y = Inches(4.0)
for line in info_lines:
    add_textbox(
        slide,
        Inches(3.5),
        y,
        Inches(6.3),
        Inches(0.35),
        line,
        size=17,
        color=RGBColor(0xBB, 0xCC, 0xDD),
        align=PP_ALIGN.CENTER,
    )
    y += Inches(0.4)

# 日期
add_textbox(
    slide,
    Inches(0.8),
    Inches(7.1),
    Inches(11.7),
    Inches(0.35),
    "2026年5月",
    size=14,
    color=C_LIGHT,
    align=PP_ALIGN.CENTER,
)


# ================================================================
# 第2页：目录
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
add_section_bar(slide, "目  录", "CONTENTS")

items = [
    ("01", "项目背景与意义"),
    ("02", "系统相关技术"),
    ("03", "系统需求分析"),
    ("04", "系统总体设计"),
    ("05", "系统详细设计与实现"),
    ("06", "系统界面展示"),
    ("07", "系统测试"),
    ("08", "总结与展望"),
]
for i, (num, title) in enumerate(items):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + col * Inches(6.2)
    y = Inches(1.5) + row * Inches(1.3)
    add_rect(slide, x, y, Inches(5.8), Inches(1.0), C_LIGHT_BG, 0.08)
    add_textbox(
        slide,
        x + Inches(0.25),
        y + Inches(0.15),
        Inches(0.7),
        Inches(0.6),
        num,
        size=28,
        bold=True,
        color=C_SECONDARY,
    )
    add_textbox(
        slide,
        x + Inches(1.0),
        y + Inches(0.25),
        Inches(4.5),
        Inches(0.5),
        title,
        size=18,
        color=C_DARK,
    )


# ================================================================
# 第3页：项目背景
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
add_section_bar(slide, "一、项目背景与意义")

add_card(
    slide,
    Inches(0.6),
    Inches(1.4),
    Inches(5.8),
    Inches(5.5),
    "📌 当前研学管理面临的问题",
    [
        "研学项目信息发布渠道分散，学生难以及时获取",
        "报名审核流程依赖线下操作，效率低下且易出错",
        "研学资源（资讯、课程）缺乏统一的管理平台",
        "师生之间缺少便捷的在线交流互动渠道",
        "研学政策资讯获取渠道有限，检索不便",
    ],
    title_color=C_ACCENT,
    bg=RGBColor(0xFF, 0xF3, 0xE0),
)

add_card(
    slide,
    Inches(6.8),
    Inches(1.4),
    Inches(5.8),
    Inches(5.5),
    "🎯 平台建设目标",
    [
        "打造一站式高校研学综合服务平台",
        "项目管理全流程线上化，提升管理效率",
        "集中管理研学资讯与课程，支持在线学习",
        "论坛社区 + 实时聊天，促进师生交流互动",
        "基于 RAG 的 AI 智能问答，赋能研学知识服务",
    ],
    title_color=C_GREEN,
    bg=RGBColor(0xE8, 0xF5, 0xE9),
)


# ================================================================
# 第4页：系统相关技术
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
add_section_bar(slide, "二、系统相关技术")

techs = [
    ("后端框架", "Spring Boot 4.0.3 + MyBatis"),
    ("前端框架", "Vue 3 + TypeScript + Element Plus"),
    ("数据库", "MySQL 5.7"),
    ("缓存与实时通信", "Redis + WebSocket"),
    ("AI 服务", "Python FastAPI + ChromaDB + Ollama"),
    ("构建工具", "Maven（后端）/ Vite（前端）"),
    ("开发环境", "JDK 17 + Node.js"),
]
y = Inches(1.5)
for label, value in techs:
    add_rect(slide, Inches(0.8), y, Inches(2.8), Inches(0.55), C_SECONDARY, 0.06)
    add_textbox(
        slide,
        Inches(0.9),
        y + Inches(0.08),
        Inches(2.6),
        Inches(0.4),
        label,
        size=15,
        bold=True,
        color=C_WHITE,
        align=PP_ALIGN.CENTER,
    )
    add_rect(slide, Inches(3.7), y, Inches(8.8), Inches(0.55), C_LIGHT_BG, 0.06)
    add_textbox(
        slide,
        Inches(3.8),
        y + Inches(0.08),
        Inches(8.6),
        Inches(0.4),
        value,
        size=15,
        color=C_DARK,
    )
    y += Inches(0.62)

# 技术架构说明
add_textbox(
    slide,
    Inches(0.8),
    Inches(6.0),
    Inches(11.5),
    Inches(0.8),
    "系统采用前后端分离架构：前端 Vue 3 SPA 通过 RESTful API 与后端通信，后端 Spring Boot 提供业务服务，\n"
    "AI 智能问答模块作为独立微服务运行，数据库采用 MySQL，缓存与在线状态使用 Redis，实时通信基于 WebSocket。",
    size=14,
    color=C_MED,
)


# ================================================================
# 第5页：系统需求分析
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
add_section_bar(slide, "三、系统需求分析")

# 三角色
roles = [
    ("管理员", "用户管理、项目审核\n资讯/课程管理\n数据统计大盘", C_SECONDARY),
    ("教师", "创建管理研学项目\n发布资讯、上传课程\n参与论坛讨论", C_TEAL),
    ("学生", "浏览项目并报名\n查看资讯和课程\n论坛交流、AI问答、聊天", C_ACCENT),
]
for i, (role, desc, color) in enumerate(roles):
    x = Inches(0.6) + i * Inches(4.2)
    add_rect(slide, x, Inches(1.4), Inches(3.9), Inches(2.6), color, 0.08)
    add_textbox(
        slide,
        x + Inches(0.2),
        Inches(1.55),
        Inches(3.5),
        Inches(0.45),
        role,
        size=24,
        bold=True,
        color=C_WHITE,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        x + Inches(0.3),
        Inches(2.15),
        Inches(3.3),
        Inches(1.6),
        desc,
        size=15,
        color=C_WHITE,
        align=PP_ALIGN.CENTER,
    )

# 功能模块
add_textbox(
    slide,
    Inches(0.6),
    Inches(4.3),
    Inches(12),
    Inches(0.4),
    "系统功能模块",
    size=20,
    bold=True,
    color=C_PRIMARY,
)
modules = [
    "研学论坛",
    "消息通知",
    "研学资讯",
    "精品课程",
    "AI智能问答",
    "研学项目",
    "系统管理",
]
for i, m in enumerate(modules):
    x = Inches(0.6) + i * Inches(1.78)
    add_rect(slide, x, Inches(4.85), Inches(1.65), Inches(0.55), C_LIGHT_BG, 0.08)
    add_textbox(
        slide,
        x,
        Inches(4.92),
        Inches(1.65),
        Inches(0.4),
        m,
        size=14,
        bold=True,
        color=C_SECONDARY,
        align=PP_ALIGN.CENTER,
    )

add_textbox(
    slide,
    Inches(0.6),
    Inches(5.6),
    Inches(12),
    Inches(0.8),
    "系统支持三种角色：管理员拥有全部管理权限，可进行用户、项目、资讯、课程的管理和审核工作；\n"
    "教师可创建和管理研学项目、发布资讯与课程、参与论坛讨论；学生可浏览资源、报名项目、交流互动和使用 AI 问答。",
    size=14,
    color=C_MED,
)


# ================================================================
# 第6页：系统架构图（重制版）
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
add_section_bar(slide, "四、系统总体设计 — 系统架构")

# 前端层
add_rect(slide, Inches(1.5), Inches(1.2), Inches(10.3), Inches(1.6), C_SECONDARY, 0.08)
add_textbox(
    slide,
    Inches(1.7),
    Inches(1.3),
    Inches(10),
    Inches(0.4),
    "前端展示层（Vue 3 + TypeScript + Element Plus）",
    size=18,
    bold=True,
    color=C_WHITE,
)
add_textbox(
    slide,
    Inches(1.7),
    Inches(1.7),
    Inches(10),
    Inches(0.9),
    "管理员布局（AdminLayout）      教师布局（TeacherLayout）      学生布局（StudentLayout）\n"
    "路由权限控制（Vue Router 导航守卫）       HTTP 请求拦截器（自动携带 JWT Token）",
    size=14,
    color=C_WHITE,
)

# 箭头1
add_textbox(
    slide,
    Inches(5.5),
    Inches(2.8),
    Inches(2.3),
    Inches(0.4),
    "⬇  RESTful API  ⬇",
    size=14,
    color=C_MED,
    align=PP_ALIGN.CENTER,
)

# 后端层
add_rect(slide, Inches(1.5), Inches(3.2), Inches(6.5), Inches(1.8), C_TEAL, 0.08)
add_textbox(
    slide,
    Inches(1.7),
    Inches(3.3),
    Inches(6.0),
    Inches(0.4),
    "后端服务层（Spring Boot 4.0.3 + MyBatis）",
    size=18,
    bold=True,
    color=C_WHITE,
)
add_textbox(
    slide,
    Inches(1.7),
    Inches(3.7),
    Inches(6.0),
    Inches(1.2),
    "三层架构：Controller → Service → Mapper\n"
    "公共字段 AOP 自动填充（createTime / updateTime）\n"
    "JWT 令牌拦截器统一鉴权（/admin/** /chat/**）\n"
    "Redis 缓存（在线状态追踪，TTL=15分钟）",
    size=14,
    color=C_WHITE,
)

# AI微服务
add_rect(slide, Inches(8.3), Inches(3.2), Inches(3.5), Inches(1.8), C_PURPLE, 0.08)
add_textbox(
    slide,
    Inches(8.5),
    Inches(3.3),
    Inches(3.1),
    Inches(0.4),
    "AI 智能问答微服务",
    size=16,
    bold=True,
    color=C_WHITE,
    align=PP_ALIGN.CENTER,
)
add_textbox(
    slide,
    Inches(8.5),
    Inches(3.7),
    Inches(3.1),
    Inches(1.1),
    "Python FastAPI\n" "ChromaDB 向量检索\n" "Ollama qwen3 大模型\n" "SSE 流式输出",
    size=13,
    color=C_WHITE,
    align=PP_ALIGN.CENTER,
)

# 箭头2
add_textbox(
    slide,
    Inches(4.0),
    Inches(5.0),
    Inches(2.3),
    Inches(0.4),
    "⬇  数据持久化  ⬇",
    size=14,
    color=C_MED,
    align=PP_ALIGN.CENTER,
)

# 数据层
add_rect(slide, Inches(1.5), Inches(5.4), Inches(10.3), Inches(1.5), C_ACCENT, 0.08)
add_textbox(
    slide,
    Inches(1.7),
    Inches(5.5),
    Inches(10),
    Inches(0.4),
    "数据持久层（MySQL + Redis）",
    size=18,
    bold=True,
    color=C_WHITE,
)
add_textbox(
    slide,
    Inches(1.7),
    Inches(5.9),
    Inches(10),
    Inches(0.9),
    "核心数据表：user · project_study · project_signup · research_news · premium_course\n"
    "社交数据表：forum_post · forum_comment · friend_relation · private_message\n"
    "所有表使用 InnoDB 引擎 + utf8mb4 字符集，create_time 由 AOP 切面自动填充",
    size=14,
    color=C_WHITE,
)

# 数据流说明
add_textbox(
    slide,
    Inches(0.6),
    Inches(7.0),
    Inches(12),
    Inches(0.35),
    "数据流：前端 HTTP 请求 → JWT 拦截器鉴权 → Controller 接收参数 → Service 业务处理 → Mapper 数据库操作 → 响应返回前端",
    size=13,
    color=C_MED,
    align=PP_ALIGN.CENTER,
)


# ================================================================
# 第7页：数据库设计
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
add_section_bar(
    slide, "四、系统总体设计 — 数据库设计", "数据库名称：study_flow，共 9 张核心业务表"
)

tables = [
    ("user", "用户表", "存储管理员/教师/学生三种角色的账号信息和个人资料"),
    (
        "project_study",
        "研学项目表",
        "存储研学项目的主题、类别、正文、状态、点赞数、点击量等信息",
    ),
    (
        "project_signup",
        "项目报名表",
        "记录学生报名信息，关联用户与项目，管理报名审核状态",
    ),
    ("research_news", "研学资讯表", "存储资讯内容，支持分类管理和审核发布流程"),
    (
        "premium_course",
        "精品课程表",
        "存储课程资源信息，含视频路径、讲师、时长等专属字段",
    ),
    ("forum_post", "论坛帖子表", "存储帖子标题、正文、分类（8种）和互动统计数据"),
    ("forum_comment", "论坛评论表", "存储评论内容，关联帖子和评论作者"),
    ("friend_relation", "好友关系表", "管理好友请求、通过和拒绝的状态流转"),
    ("private_message", "私信消息表", "存储用户间聊天消息，支持已读/未读标记"),
]

y = Inches(1.4)
for i, (name, cname, desc) in enumerate(tables):
    bg = C_LIGHT_BG if i % 2 == 0 else C_WHITE
    add_rect(slide, Inches(0.6), y, Inches(12), Inches(0.52), bg, 0.03)
    add_textbox(
        slide,
        Inches(0.8),
        y + Inches(0.06),
        Inches(2.5),
        Inches(0.4),
        name,
        size=14,
        bold=True,
        color=C_SECONDARY,
    )
    add_textbox(
        slide,
        Inches(3.5),
        y + Inches(0.06),
        Inches(2.5),
        Inches(0.4),
        cname,
        size=14,
        color=C_DARK,
    )
    add_textbox(
        slide,
        Inches(6.2),
        y + Inches(0.06),
        Inches(6.2),
        Inches(0.4),
        desc,
        size=13,
        color=C_MED,
    )
    y += Inches(0.54)

add_textbox(
    slide,
    Inches(0.6),
    y + Inches(0.15),
    Inches(12),
    Inches(0.35),
    "所有表均使用 InnoDB 存储引擎与 utf8mb4 字符集，以 user 表为核心的辐射状关联结构",
    size=13,
    color=C_MED,
)


# ================================================================
# 第8页：功能模块 — 研学论坛 & 消息通知
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
add_section_bar(slide, "五、系统详细设计与实现（1/3）")

add_card(
    slide,
    Inches(0.6),
    Inches(1.3),
    Inches(5.9),
    Inches(5.7),
    "💬 研学论坛模块",
    [
        "帖子发布：支持 8 种分类（竞赛心得、实训心得、读书心得等）",
        "卡片式展示：CSS Grid 网格布局，14 种柔和色调随机分配",
        "2-3 度随机旋转偏移，营造仿手工便签视觉效果",
        "支持一键切换倾斜效果的开关按钮",
        "帖子详情通过右侧抽屉展示正文和评论列表",
        "评论按发表时间正序排列",
        "级联删除：删除帖子时同步清理所有关联评论",
    ],
    title_color=C_SECONDARY,
)

add_card(
    slide,
    Inches(6.8),
    Inches(1.3),
    Inches(5.9),
    Inches(5.7),
    "✉️ 消息通知模块",
    [
        "好友关系管理：添加好友、通过/拒绝请求、删除好友",
        "基于 WebSocket 的一对一实时消息通信",
        "JWT 令牌握手认证，保障连接安全性",
        "消息先持久化存储再实时推送，确保不丢失",
        "离线消息保存：接收方不在线时消息暂存数据库",
        "断线自动重连：意外断开时自动递归恢复连接",
        "Redi s 缓存在线状态（TTL=15分钟），双重触发更新",
    ],
    title_color=C_TEAL,
)


# ================================================================
# 第9页：功能模块 — 资讯 & 课程 & AI
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
add_section_bar(slide, "五、系统详细设计与实现（2/3）")

# 资讯
add_rect(slide, Inches(0.6), Inches(1.3), Inches(3.8), Inches(2.6), C_LIGHT_BG, 0.06)
add_textbox(
    slide,
    Inches(0.8),
    Inches(1.42),
    Inches(3.4),
    Inches(0.35),
    "📰 研学资讯模块",
    size=17,
    bold=True,
    color=C_SECONDARY,
)
add_multi_text(
    slide,
    Inches(0.9),
    Inches(1.85),
    Inches(3.3),
    Inches(1.9),
    [
        "• 富文本在线编辑资讯内容",
        "• 5 种预定义分类 + 自定义扩展",
        "• 状态流转：待审核→已发布→已下架",
        "• 点击量自动统计",
        "• 排序权重 + 更新时间排序",
    ],
    size=13,
    color=C_DARK,
)

# 课程
add_rect(
    slide,
    Inches(4.7),
    Inches(1.3),
    Inches(3.8),
    Inches(2.6),
    RGBColor(0xE8, 0xF5, 0xE9),
    0.06,
)
add_textbox(
    slide,
    Inches(4.9),
    Inches(1.42),
    Inches(3.4),
    Inches(0.35),
    "🎓 精品课程模块",
    size=17,
    bold=True,
    color=C_GREEN,
)
add_multi_text(
    slide,
    Inches(5.0),
    Inches(1.85),
    Inches(3.3),
    Inches(1.9),
    [
        "• 封面图片 + 课程视频上传",
        "• 讲师姓名 + 课程时长标注",
        "• HTML5 视频在线播放",
        "• 状态：草稿→待审核→已发布→已下架",
        "• 统一文件上传接口存储资源",
    ],
    size=13,
    color=C_DARK,
)

# AI
add_rect(
    slide,
    Inches(8.8),
    Inches(1.3),
    Inches(3.8),
    Inches(2.6),
    RGBColor(0xF3, 0xE5, 0xF5),
    0.06,
)
add_textbox(
    slide,
    Inches(9.0),
    Inches(1.42),
    Inches(3.4),
    Inches(0.35),
    "🤖 AI 智能问答模块",
    size=17,
    bold=True,
    color=C_PURPLE,
)
add_multi_text(
    slide,
    Inches(9.1),
    Inches(1.85),
    Inches(3.3),
    Inches(1.9),
    [
        "• RAG 检索增强生成架构",
        "• ChromaDB 向量检索（余弦相似度）",
        "• Ollama qwen3:8b 大模型生成",
        "• SSE 流式输出逐字推送",
        "• 支持普通 + 流式两种模式",
    ],
    size=13,
    color=C_DARK,
)

# AI 流程
add_rect(slide, Inches(0.6), Inches(4.2), Inches(12), Inches(2.8), C_LIGHT_BG, 0.06)
add_textbox(
    slide,
    Inches(0.8),
    Inches(4.35),
    Inches(11.6),
    Inches(0.35),
    "AI 智能问答流程",
    size=17,
    bold=True,
    color=C_PURPLE,
)

flow_steps = [
    ("用户提问", "前端输入\n问题文本"),
    ("向量化", "qwen3-embed:4b\n转为向量"),
    ("向量检索", "ChromaDB\nTop-K=5"),
    ("提示词组装", "知识+问题\n拼接为 prompt"),
    ("大模型生成", "qwen3:8b\n生成回答"),
    ("结果返回", "普通/流式\n输出给前端"),
]
for i, (step, desc) in enumerate(flow_steps):
    x = Inches(0.8) + i * Inches(2.05)
    add_rect(
        slide,
        x,
        Inches(4.85),
        Inches(1.85),
        Inches(1.4),
        C_SECONDARY if i != 4 else C_ACCENT,
        0.08,
    )
    add_textbox(
        slide,
        x + Inches(0.1),
        Inches(4.95),
        Inches(1.65),
        Inches(0.35),
        step,
        size=14,
        bold=True,
        color=C_WHITE,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        x + Inches(0.1),
        Inches(5.35),
        Inches(1.65),
        Inches(0.8),
        desc,
        size=12,
        color=C_WHITE,
        align=PP_ALIGN.CENTER,
    )
    if i < len(flow_steps) - 1:
        add_textbox(
            slide,
            x + Inches(1.8),
            Inches(5.3),
            Inches(0.3),
            Inches(0.3),
            "→",
            size=18,
            bold=True,
            color=C_MED,
        )

add_textbox(
    slide,
    Inches(0.8),
    Inches(6.35),
    Inches(11.6),
    Inches(0.35),
    "知识来源：爬虫采集桂电快讯 + 人民研学网    文本分块：700字符/块，重叠120字符    向量库：studyflow_knowledge 集合",
    size=12,
    color=C_MED,
    align=PP_ALIGN.CENTER,
)


# ================================================================
# 第10页：功能模块 — 研学项目 & 系统管理
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
add_section_bar(slide, "五、系统详细设计与实现（3/3）")

add_card(
    slide,
    Inches(0.6),
    Inches(1.3),
    Inches(5.9),
    Inches(5.7),
    "📋 研学项目管理模块（核心业务）",
    [
        "教师/管理员创建项目：填写主题、类别、简介、富文本正文和封面",
        "新建项目默认状态为「待审核」，管理员审核通过后变为「已发布」",
        "项目列表支持按关键词和类别筛选，分页展示",
        "排序规则：待审核报名优先 → 点击量降序 → 点赞数降序",
        "学生报名：后端校验防重复，生成「待审核」报名记录",
        "报名审核：管理员通过（状态更新）或拒绝（删除记录）",
        "互动功能：点赞 +1、点击量自动递增、项目分享到好友",
    ],
    title_color=C_SECONDARY,
)

add_card(
    slide,
    Inches(6.8),
    Inches(1.3),
    Inches(5.9),
    Inches(5.7),
    "⚙️ 系统管理模块",
    [
        "用户管理：分页查询，多条件筛选（用户名/姓名/性别/状态）",
        "用户编辑、账号启用/禁用、删除操作",
        "研学项目管理：审核通过、删除、报名用户管理",
        "研学资讯管理：新增、编辑、审核、下架、删除",
        "精品课程管理：新增、编辑、审核、上下架、删除",
        "前端路由守卫 + 后端 JWT 拦截器双重鉴权",
        "首页 ECharts 数据统计大盘：项目数/在线/点击/参与",
    ],
    title_color=C_ACCENT,
)


# ================================================================
# 第11页：系统界面展示
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
add_section_bar(slide, "六、系统界面展示", "以下为系统各主要功能页面的运行截图")

imgs = [
    ("登录注册页面", "用户登录与注册界面"),
    ("系统首页", "数据统计大盘 + ECharts图表"),
    ("研学项目列表", "项目卡片网格展示"),
    ("项目详情页", "富文本内容 + 报名/点赞"),
    ("论坛社区", "卡片式帖子列表"),
    ("好友聊天", "WebSocket 实时通信"),
    ("精品课程", "课程列表 + 视频播放"),
    ("AI 智能问答", "RAG 对话界面"),
]
for i, (title, desc) in enumerate(imgs):
    col = i % 4
    row = i // 4
    x = Inches(0.5) + col * Inches(3.15)
    y = Inches(1.4) + row * Inches(2.8)
    add_rect(slide, x, y, Inches(2.9), Inches(2.5), C_LIGHT_BG, 0.08)
    add_textbox(
        slide,
        x + Inches(0.15),
        y + Inches(0.8),
        Inches(2.6),
        Inches(0.35),
        f"[截图] {title}",
        size=14,
        bold=True,
        color=C_SECONDARY,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        x + Inches(0.15),
        y + Inches(1.2),
        Inches(2.6),
        Inches(0.35),
        desc,
        size=12,
        color=C_MED,
        align=PP_ALIGN.CENTER,
    )

add_textbox(
    slide,
    Inches(0.8),
    Inches(7.0),
    Inches(11.7),
    Inches(0.35),
    "（此处放置系统实际运行截图，展示平台各功能模块的界面效果）",
    size=14,
    color=C_MED,
    align=PP_ALIGN.CENTER,
)


# ================================================================
# 第12页：系统测试
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
add_section_bar(slide, "七、系统测试")

# 统计卡片
stats = [
    ("9 个", "测试模块"),
    ("52 个", "测试用例"),
    ("4 种", "测试方法"),
    ("100%", "通过率"),
]
for i, (num, label) in enumerate(stats):
    x = Inches(0.8) + i * Inches(3.1)
    add_rect(slide, x, Inches(1.3), Inches(2.8), Inches(1.5), C_SECONDARY, 0.1)
    add_textbox(
        slide,
        x,
        Inches(1.45),
        Inches(2.8),
        Inches(0.6),
        num,
        size=38,
        bold=True,
        color=C_WHITE,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        x,
        Inches(2.1),
        Inches(2.8),
        Inches(0.4),
        label,
        size=16,
        color=RGBColor(0xCC, 0xDD, 0xFF),
        align=PP_ALIGN.CENTER,
    )

# 测试模块
add_textbox(
    slide,
    Inches(0.8),
    Inches(3.1),
    Inches(11.5),
    Inches(0.35),
    "测试覆盖模块",
    size=18,
    bold=True,
    color=C_PRIMARY,
)
mods = [
    "用户系统",
    "研学项目",
    "论坛社区",
    "消息通知",
    "校园资讯",
    "精品课程",
    "AI问答",
    "爬虫采集",
    "系统管理",
]
y = Inches(3.6)
for i, m in enumerate(mods):
    col = i % 5
    row = i // 5
    x = Inches(0.8) + col * Inches(2.4)
    yy = y + row * Inches(0.5)
    add_rect(slide, x, yy, Inches(2.2), Inches(0.42), C_LIGHT_BG, 0.06)
    add_textbox(
        slide,
        x,
        yy + Inches(0.04),
        Inches(2.2),
        Inches(0.35),
        m,
        size=14,
        color=C_SECONDARY,
        align=PP_ALIGN.CENTER,
    )

# 测试方法
add_textbox(
    slide,
    Inches(0.8),
    Inches(4.8),
    Inches(11.5),
    Inches(0.35),
    "测试方法",
    size=18,
    bold=True,
    color=C_PRIMARY,
)
add_multi_text(
    slide,
    Inches(0.8),
    Inches(5.2),
    Inches(11.5),
    Inches(1.5),
    [
        "• 功能测试：针对每个模块的核心业务逻辑设计测试用例，验证完整链路的正确性",
        "• 接口测试：验证 RESTful API 在合法和非法参数下的响应符合规范",
        "• 集成测试：验证前后端数据交互和各模块间协同工作的正确性",
        "• 兼容性测试：在 Chrome 浏览器下验证页面布局和交互功能的完整性",
    ],
    size=14,
    color=C_DARK,
)

add_textbox(
    slide,
    Inches(0.8),
    Inches(6.8),
    Inches(11.5),
    Inches(0.35),
    "测试结论：所有 52 个测试用例全部通过，系统功能完整、运行稳定，满足设计要求。",
    size=15,
    bold=True,
    color=C_GREEN,
)


# ================================================================
# 第13页：总结与展望
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
add_section_bar(slide, "八、总结与展望")

add_card(
    slide,
    Inches(0.6),
    Inches(1.3),
    Inches(5.9),
    Inches(5.7),
    "✅ 已完成工作",
    [
        "完成了研学论坛、消息通知、研学资讯、精品课程、AI智能问答、研学项目、系统管理七大模块",
        "实现了前后端分离架构和 RESTful API 接口设计",
        "集成了基于 RAG 的 AI 智能问答微服务（ChromaDB + Ollama）",
        "实现了 WebSocket 实时消息通信和好友关系管理",
        "构建了完整的 MySQL 数据库表结构（9 张核心业务表）",
        "通过了 52 个测试用例的系统功能测试，通过率 100%",
        "实现了 JWT 鉴权 + AOP 自动填充等公共基础设施",
    ],
    title_color=C_GREEN,
    bg=RGBColor(0xE8, 0xF5, 0xE9),
)

add_card(
    slide,
    Inches(6.8),
    Inches(1.3),
    Inches(5.9),
    Inches(5.7),
    "🔭 未来展望",
    [
        "支持小组研学项目协作与任务分配功能",
        "引入移动端小程序或 APP 适配，拓展使用场景",
        "优化 AI 问答的检索精度和响应速度",
        "增加研学成果在线提交与评价功能",
        "引入数据可视化大屏展示平台运营数据",
        "支持多院校多租户扩展部署模式",
    ],
    title_color=C_SECONDARY,
)


# ================================================================
# 第14页：致谢
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_PRIMARY)

# 装饰
add_rect(
    slide,
    Inches(0),
    Inches(0),
    Inches(0.2),
    prs.slide_height,
    RGBColor(0x15, 0x30, 0x5A),
)
add_rect(slide, Inches(0), Inches(3.3), prs.slide_width, Inches(0.04), C_ACCENT)
add_rect(
    slide,
    Inches(0),
    Inches(7.0),
    prs.slide_width,
    Inches(0.5),
    RGBColor(0x15, 0x30, 0x5A),
)

add_textbox(
    slide,
    Inches(1.5),
    Inches(2.0),
    Inches(10.3),
    Inches(0.8),
    "感谢各位老师聆听！",
    size=44,
    bold=True,
    color=C_WHITE,
    align=PP_ALIGN.CENTER,
)
add_textbox(
    slide,
    Inches(1.5),
    Inches(3.7),
    Inches(10.3),
    Inches(0.5),
    "请各位老师批评指正",
    size=22,
    color=RGBColor(0xBB, 0xCC, 0xEE),
    align=PP_ALIGN.CENTER,
)
add_textbox(
    slide,
    Inches(1.5),
    Inches(5.0),
    Inches(10.3),
    Inches(0.4),
    "StudyFlow 高校研学交流平台",
    size=16,
    color=C_LIGHT,
    align=PP_ALIGN.CENTER,
)
add_textbox(
    slide,
    Inches(1.5),
    Inches(5.4),
    Inches(10.3),
    Inches(0.4),
    "桂林电子科技大学",
    size=16,
    color=C_LIGHT,
    align=PP_ALIGN.CENTER,
)


# ── 保存 ──
output = r"c:\D\Code\IDEA_project\StudyFlow\StudyFlow_答辩PPT.pptx"
prs.save(output)
print(f"[OK] PPT generated: {output}")
print(f"[OK] Total {len(prs.slides)} slides")
